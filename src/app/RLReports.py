#########imports###############################################################
import colorsys
import cv2
import exifread
import glob
import imageio
import os
from operator import truediv
from PIL import Image
import pptx
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT as PP_ALIGN
import pptx.util
from pptx.util import Inches, Pt
from pylab import *
import re
from scipy.signal import find_peaks
import shutil
from datetime import datetime

from src import constants
from src.logs import logger
from src.gui.utils_gui import send_progress_to_gui, send_error_to_gui

np.warnings.filterwarnings('ignore')


###############################################################################
###############################################################################
# Example data:
# report_config = {
#     "config": {
#         "image_path": r"C:\Users\mms00519\Downloads\Indoor_cases",
#         "thumbnail_path": r"C:\Users\mms00519\Downloads\Indoor_cases\Thumbnail",
#         "presentation_name": "RLT_presentation",
#         "attribute_on": 1,
#         "avg_luma": True,
#         "contrast": True,
#         "black_level": True,
#         "white_level": True,
#         "over_exposed": True,
#         "under_exposed": True,
#         "dynamic_range": False,
#         "peak_saturation1": False,
#         "peak_hue1": False,
#         "peak_saturation2": False,
#         "peak_hue2": False,
#         "sharpness": True,
#         "ISO": True,
#         "ET": True
#     },
#     "summary_params": {
#         "af": True,
#         "ae": True,
#         "awb": True,
#         "colors": True,
#         "noise": True,
#         "details": True,
#         "artifacts": True,
#         "torch": True,
#         "flash": True
#     },
#     "summary_items": {
#         "attribute": True,
#         "level": True,
#         "issues": True,
#         "suggestions": True
#     },
#     "attribute": {
#         "exposure": True,
#         "colors": True,
#         "noise": True,
#         "details": True,
#         "artifacts": True
#     }
# }
#
def _generate_rlt_report(report_config: dict, gui_window=None, gui_event=None):
    logger.info("Creating object... ")

    new_report = RLReports(report_config, gui_window, gui_event)

    logger.info("Creating Presentation... ")
    new_report.create_presentation()

    del new_report


def generate_rlt_report(report_config: dict, gui_window=None, gui_event=None):
    rlt_thread = threading.Thread(target=_generate_rlt_report, args=(report_config, gui_window, gui_event), daemon=True)
    rlt_thread.name = 'RLTReportsGeneration'

    logger.info(f"Starting {rlt_thread.name} Thread")
    rlt_thread.start()

    # TODO: Kill thread after it is finished


# RLTReport Class
class RLReports:
    # Class initialization
    def __init__(self, config_dict, gui_window=None, gui_event=None):
        self.gui_window = gui_window
        self.gui_event = gui_event

        try:
            self.config = config_dict['config']
            self.summary_params = config_dict['summary_params']
            self.summary_items = config_dict['summary_items']
            self.attribute = config_dict['attribute']
        except KeyError:
            logger.critical("RLReports class got wrong data!")
            raise ValueError("RLReports class got wrong data in dict!")

    # Create presentation
    def create_presentation(self):

        if self.gui_window is not None and self.gui_event is not None:
            logger.info("Starting RLReports Generation")
            send_progress_to_gui(self.gui_window, self.gui_event, 0, 'Starting')

        image_path = self.config['image_path']
        image_outpath = self.config['thumbnail_path']
        presentation_name = self.config['presentation_name']
        if presentation_name == "":
            presentation_name = "RealLife Report"

        date_str = datetime.now().strftime("%Y%m%d-%H%M%S")
        output_file = os.path.join(self.config['image_path'], os.path.pardir, f"{presentation_name}_{date_str}.pptm")

        prs = pptx.Presentation(os.path.join(constants.VENDOR_DIR, 'rltreport', 'RLTv1.pptm'))

        prs.slide_height = Inches(7.5)
        prs.slide_width = Inches(13.33070866)

        images = [g for g in glob.glob(image_path + "/**/case*.jp*", recursive=True)]

        def atoi(input):
            return int(input) if input.isdigit() else input

        def natural_keys(input):
            return [atoi(c) for c in re.split('(\d+)', input)]

        if self.gui_window is not None and self.gui_event is not None:
            logger.info("Initializing RLReports")
            send_progress_to_gui(self.gui_window, self.gui_event, 3, 'Initializing')

        images.sort(key=natural_keys)

        slide = prs.slides.add_slide(prs.slide_layouts[7])
        title = slide.shapes.title
        title.text = "Images - Executed test cases"
        title.text_frame.paragraphs[0].font.size = Pt(26)
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

        if self.gui_window is not None and self.gui_event is not None:
            logger.info("Generating Thumbnails")
            send_progress_to_gui(self.gui_window, self.gui_event, 30, 'Generating Thumbnails')
        create_thumbnail(prs, slide, images, image_path, image_outpath)

        shutil.rmtree(image_outpath)
        if self.gui_window is not None and self.gui_event is not None:
            logger.info("Generating Report")
            send_progress_to_gui(self.gui_window, self.gui_event, 38, 'Generating Report')
        set_images_to_slide(prs, images, self)

        if self.gui_window is not None and self.gui_event is not None:
            send_progress_to_gui(self.gui_window, self.gui_event, 95, 'Saving')
        prs.save(output_file)

        if self.gui_window is not None and self.gui_event is not None:
            send_progress_to_gui(self.gui_window, self.gui_event, 100, 'Done!', 'new_file', output_file)
        logger.info(f"RLReport {presentation_name} Done\nSaved to: {output_file}")

    # Get image stats
    def show_req(self, rgb_copy, img_iso):
        stats_list = []
        item_list = ["Item"]

        img = Image.open(img_iso)

        for key, value in list(self.config.items())[4:]:
            # print (item)
            if value:
                item_list.extend([str(key)])

        #### ET ISO ###########
        if self.config['ET'] or self.config['ISO']:
            with open(img_iso, "rb") as f:
                tags = exifread.process_file(f)
                iso = tags.get('EXIF ISOSpeedRatings')
                et = tags.get('EXIF ExposureTime')

        #### sharpness ########
        if self.config['Sharpness']:
            img_sharp = cv2.imread(img_iso, 0)
            sharpness = sharpness_image(img_sharp)

        #### dyn range ########
        if self.config['Dynamic Range']:
            img_dyn_range = cv2.imread(img_iso, 1)
            dyn_range = dynamic_range(img_dyn_range)
        #### over exp ########
        if self.config['Over Exposed']:
            img_dyn_range = cv2.imread(img_iso, 1)
            dyn_range = dynamic_range(img_dyn_range)
        #### under exp ########
        if self.config['Under Exposed']:
            img_dyn_range = cv2.imread(img_iso, 1)
            dyn_range = dynamic_range(img_dyn_range)

        #### BL WL contrast ###
        if self.config['Contrast'] or self.config['White Level'] or self.config['Black Level']:
            img_grey = img.convert('L')
            f1 = figure(1)
            f1.suptitle("Input Image")
            plt.imshow(img_grey, cmap='gray', vmin=0, vmax=255)
            img_array = np.asarray(img_grey)
            f2 = figure(2)
            f2.suptitle("Input Image Histogram")
            hist, bins = np.histogram(img_array, 256, [0, 256])
            plot(hist)
            f2 = figure(3)
            f2.suptitle("Input Image CDF")
            cdf = hist.cumsum()
            cdf_nom = cdf / cdf.max()
            black_level = np.min(np.argwhere(cdf_nom > 0.02))
            white_level = np.min(np.argwhere(cdf_nom > 0.98))
            contrast = int(((white_level - black_level) / 255) * 100)

        #### AVR Luma #########
        if self.config['Avg Luma']:
            im = Image.fromarray(rgb_copy.astype('uint8'))
            yuv = asarray(im.convert('L'))
            h, s, v = colorsys.rgb_to_hsv(np.mean(rgb_copy[:, :, 0]), np.mean(rgb_copy[:, :, 1]),
                                          np.mean(rgb_copy[:, :, 2]))
            avr_y = int(np.mean(yuv))

        #### saturattion ######
        if self.config['Peak Saturation 1'] \
                or self.config['Peak Saturation 2']\
                or self.config['Peak Hue 1']\
                or self.config['Peak Hue 2']:
            img_peak_sat = np.asarray(img)
            peak_sat = peak_saturation(img_peak_sat)

        if self.config['Avg Luma']:
            stats_list.extend([str(avr_y)])

        if self.config['Contrast']:
            stats_list.extend([str(contrast) + "%"])

        if self.config['Black Level']:
            stats_list.extend([str(black_level)])

        if self.config['White Level']:
            stats_list.extend([str(white_level)])

        if self.config['Dynamic Range']:
            stats_list.extend([str(dyn_range[2]) + "%"])

        if self.config['Over Exposed']:
            stats_list.extend([str(dyn_range[0]) + "%"])

        if self.config['Under Exposed']:
            stats_list.extend([str(dyn_range[1]) + "%"])

        if self.config['Peak Saturation 1']:
            stats_list.extend([str(peak_sat[0]) + "%"])

        if self.config['Peak Hue 1']:
            stats_list.extend([str(peak_sat[1]) + "°"])

        if self.config['Peak Saturation 2']:
            stats_list.extend([str(peak_sat[2]) + "%"])

        if self.config['Peak Hue 2']:
            stats_list.extend([str(peak_sat[3]) + "°"])

        if self.config['Sharpness']:
            stats_list.extend([str(sharpness)])

        if self.config['ISO']:
            stats_list.extend([str(iso)])

        if self.config['ET']:
            stats_list.extend([str(et)])

        plt.clf()
        return stats_list, item_list, self.config['attribute_on']

    # Create slides with summary
    def set_summary_to_slide(self, prs):
        # print (len(value_dict[1]))
        table_default = []
        summary_list = []
        for item in self.summary_params.keys():
            summary_list.append(item)
        for item in self.summary_items.keys():
            table_default.append(item)
        aa = len(self.summary_items.keys())
        bb = len(self.summary_params.keys())
        # slides = prs.slides
        slide = prs.slides.add_slide(prs.slide_layouts[7])
        # print ("aaaaaa", len(slides))
        title = slide.shapes.title
        title.text = "Test Summary"
        title.text_frame.paragraphs[0].font.size = Pt(26)
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        table_left = int(prs.slide_width * 0.01)
        table_top = int(prs.slide_height * 0.2)
        table_width = int(prs.slide_width / 1.02)
        table_height = int(prs.slide_height / 8)
        shape = slide.shapes.add_table((len(self.summary_params.keys()) + 1), (len(self.summary_items.keys())),
                                       table_left, table_top,
                                       table_width, table_height)
        table = shape.table
        for row in range(0, (bb + 1)):
            for cols in range(0, aa):
                fill = table.cell(row, cols).fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(210, 210, 210)
        for row in range(1):
            for cols in range(0, aa):
                cell = table.cell(row, cols)
                l = table.cell(row, cols).text_frame.paragraphs[0]
                run = l.add_run()
                run.text = table_default[cols]
                l.alignment = PP_ALIGN.CENTER
                font = run.font
                font.size = Pt(15)
                font.color.rgb = RGBColor(0, 0, 0)
        for row in range(1, (bb + 1)):
            for cols in range(1):
                cell = table.cell(row, cols)
                l = table.cell(row, cols).text_frame.paragraphs[0]
                run = l.add_run()
                run.text = summary_list[row - 1]
                l.alignment = PP_ALIGN.CENTER
                font = run.font
                font.size = Pt(15)
                font.color.rgb = RGBColor(0, 0, 0)

    # Attribute table
    def set_attribute_to_slide(self, prs, slide):
        attribute_list = []
        defaultdict = ["Attribute", "Level", "Comparison Scale"]
        count = 0
        for item in self.attribute.keys:
            attribute_list.append(item)

        table_left = int(prs.slide_width * 0.76)
        table_top = int(prs.slide_height * 0.7)
        table_width = int(prs.slide_width / 3)
        table_height = int(prs.slide_height / 10.0)

        shape = slide.shapes.add_table(len(attribute_list) + 1, 3, table_left, table_top,
                                       table_width, table_height)
        table = shape.table
        table.columns[0].width = Inches(0.9)
        table.columns[1].width = Inches(0.6)
        table.columns[2].width = Inches(1.5)

        for row in range(len(attribute_list) + 1):
            for cols in range(3):
                fill = table.cell(row, cols).fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(229, 229, 229)
                l = table.cell(row, cols).text_frame.paragraphs[0]
                run = l.add_run()
                run.text = " "
                font = run.font
                font.size = Pt(12)
                font.color.rgb = RGBColor(0, 0, 0)
        for row in range(1):
            for cols in range(3):
                cell = table.cell(row, cols)
                l = table.cell(row, cols).text_frame.paragraphs[0]
                run = l.add_run()
                run.text = defaultdict[count]
                font = run.font
                font.size = Pt(12)
                font.color.rgb = RGBColor(0, 0, 0)
                count += 1
        count = 0
        for row in range(1, len(attribute_list) + 1):
            for cols in range(1):
                cell = table.cell(row, cols)
                l = table.cell(row, cols).text_frame.paragraphs[0]
                run = l.add_run()
                run.text = attribute_list[count]
                font = run.font
                font.size = Pt(12)
                font.color.rgb = RGBColor(0, 0, 0)
                count += 1
        count = 0


##### Static Methods #####
# Create xlsx with data

# Create thumbnails
def create_thumbnail(prs, slide, files, path, outpath):
    image_number = 0
    picture_to_slide = 1
    num_file = 0

    logger.info("Generating thumbnails...")

    for g in files:
        image_number += 1
        if os.path.exists(outpath) is False:
            os.mkdir(outpath)

        img = Image.open(g)
        head, tail = os.path.split(g)
        img_ratio = imageio.imread(g)
        if img_ratio.shape[0] > img_ratio.shape[1]:
            if (img_ratio.shape[0] / img_ratio.shape[1]) < 0.6:
                resize_width = 1080
                resize_height = 1440
            else:
                resize_width = 1080
                resize_height = 1920
        else:
            if (img_ratio.shape[0] / img_ratio.shape[1]) < 0.6:
                resize_width = 1440
                resize_height = 1080
            else:
                resize_width = 1920
                resize_height = 1080
        img1 = img.resize((resize_width, resize_height), Image.ANTIALIAS)
        img1.save(outpath + str(image_number) + tail)

    thumbnails = [g for g in glob.glob(outpath + "/**/*case*.jpg", recursive=True)]

    def atoi(thumbnails):
        return int(thumbnails) if thumbnails.isdigit() else thumbnails

    def natural_keys(thumbnails):
        return [atoi(c) for c in re.split('(\d+)', thumbnails)]

    thumbnails.sort(key=natural_keys)

    pic_left = int(prs.slide_width * 0.01)
    pic_top = int(prs.slide_height * 0.08)
    pic_width = int(prs.slide_width / picture_to_slide)
    pic_height = int(pic_width / 1.33)
    pic_separator = int(prs.slide_width * 0.003)
    pic_coeff = float(1.9 + (image_number / 100.0))
    while ((((pic_height * pic_width) * image_number / 2) + (image_number * pic_separator)
            + (pic_left * prs.slide_height) + (pic_top * prs.slide_width))
           > (prs.slide_height * prs.slide_width)) or \
            (((((image_number / 2) / (prs.slide_width / pic_width)) * pic_height)
              + pic_coeff * pic_height) > prs.slide_height):
        picture_to_slide += 0.01
        pic_width = int(prs.slide_width / picture_to_slide)
        pic_height = int(pic_width / 1.33)

    for g in thumbnails:
        num_file += 1
        img = Image.open(g)
        head, tail = os.path.split(g)
        img_ratio = imageio.imread(g)

        if img_ratio.shape[0] < img_ratio.shape[1]:
            if (img_ratio.shape[0] / img_ratio.shape[1]) < 0.6:
                pic_width = int(prs.slide_width / picture_to_slide)
                pic_height = int(pic_width / 1.33)
            else:
                pic_width = int(prs.slide_width / picture_to_slide)
                pic_height = int(pic_width / 1.78)
        else:
            if (img_ratio.shape[0] / img_ratio.shape[1]) < 0.6:
                pic_width = int(prs.slide_width / (picture_to_slide * 1.78))
                pic_height = int(pic_width * 1.33)
            else:
                pic_width = int(prs.slide_width / (picture_to_slide * 1.78))
                pic_height = int(pic_width * 1.78)
        if (num_file % 2) == 1:
            pic = slide.shapes.add_picture(g, pic_left,
                                           pic_top, pic_width, pic_height)
            pic_left = int(pic_left + pic_width + pic_separator)
            if pic_left > (prs.slide_width - pic_width):
                pic_top = int(pic_top + pic_height + pic_separator)
                pic_left = int(prs.slide_width * 0.01)


# Create slides with images
def set_images_to_slide(prs, files, obj):
    string = "case"
    pic_left = int(prs.slide_width * 0.01)
    pic_top = int(prs.slide_height * 0.08)
    slide_counter = 1
    inserted_images = 0
    image_on_slide = True
    device1_name = ""
    device2_name = ""

    progress = 40

    # get_item()
    first_stats_list = []
    second_stats_list = []
    check_summary_slide = ""
    check_summary_flag = ""
    # set_summary_to_slide(prs)
    fff = False

    logger.info("Generating slides...")

    while True:
        slide = prs.slides.add_slide(prs.slide_layouts[7])
        title = slide.shapes.title

        table_left = int(prs.slide_width * 0.01)
        table_top = int(prs.slide_height * 0.81)
        table_width = int(prs.slide_width / 1.42)
        table_height = int(prs.slide_height / 8)
        shape = slide.shapes.add_table(1, 2, table_left, table_top,
                                       table_width, table_height)
        table = shape.table
        fill = table.cell(0, 0).fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(221, 221, 221)
        fill = table.cell(0, 1).fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(221, 221, 221)

        txBox = slide.shapes.add_textbox(int(prs.slide_width * 0.01), int(prs.slide_height * 0.76),
                                         int(prs.slide_width / 4), int(prs.slide_height / 12))
        tf = txBox.text_frame
        tf.text = "Pros"
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT

        txBox = slide.shapes.add_textbox(int(prs.slide_width * 0.47), int(prs.slide_height * 0.76),
                                         int(prs.slide_width / 4), int(prs.slide_height / 12))
        tf = txBox.text_frame
        tf.text = "Cons"
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT

        progress_step = 50/len(files)

        for g in files:
            img = Image.open(g)
            head, tail = os.path.split(g)
            if not tail.find(string + str(slide_counter) + '_'):
                # print ("#####", check_summary_slide)
                # print ("#####", head)
                if check_summary_slide != head:
                    # print ("ffsssssssssssssssssssssssssssssssssssssssssssssssssssssssssssssssssssssssss")
                    fff = True
                    check_summary_slide = head
                    obj.set_summary_to_slide(prs)
                else:
                    fff = False
                    # print ("pffffffffffffffffffffffffffffffffffffffffffffffffffffffffffffffffffffffffff")
                # sprint (head)
                set_title_to_slide(tail, title, g, prs)
                width, height = img.size
                set_device_name(prs, tail, slide, image_on_slide, width, height)

                if height < width:
                    if (height / width) > 0.65:
                        pic_width = int(prs.slide_width / 3)
                        pic_height = int(pic_width * height / width)
                    else:
                        pic_width = int(prs.slide_width / 2.75)
                        pic_top = int(prs.slide_height * 0.115)
                        pic_height = int(pic_width * height / width)  # 16:9

                    if not image_on_slide:
                        if (height / width) > 0.65:
                            pic_left += int(prs.slide_width * 0.34)
                        else:
                            pic_top = int(prs.slide_height * 0.115)
                            pic_left += int(prs.slide_width * 0.37)  # 16:9
                        second_stats_list = set_image_stats(g, obj)
                        # xlsx_data(second_stats_list, tail, image_on_slide, slide_counter)
                        logger.info(f"image: {tail}")
                        device2_name = get_device_name(tail)
                    else:
                        pic_left = int(prs.slide_width * 0.01)
                        first_stats_list = set_image_stats(g, obj)
                        # xlsx_data(first_stats_list, tail, image_on_slide, slide_counter)
                        logger.info(f"image: {tail}")
                        device1_name = get_device_name(tail)
                else:
                    pic_width = int(prs.slide_width / 5.3)
                    pic_height = int(pic_width / (width / height))
                    if not image_on_slide:
                        pic_left += int(prs.slide_width * 0.195)
                        second_stats_list = set_image_stats(g, obj)
                        # xlsx_data(second_stats_list, tail, image_on_slide, slide_counter)
                        logger.info(f"image: {tail}")
                        device2_name = get_device_name(tail)
                    else:
                        pic_left = int(prs.slide_width * 0.01)
                        first_stats_list = set_image_stats(g, obj)
                        # xlsx_data(first_stats_list, tail, image_on_slide, slide_counter)
                        logger.info(f"image: {tail}")
                        device1_name = get_device_name(tail)

                progress += progress_step
                if obj.gui_window is not None and obj.gui_event is not None:
                    send_progress_to_gui(obj.gui_window, obj.gui_event, progress, f"img: {tail}")

                pic = slide.shapes.add_picture(g, pic_left, pic_top,
                                               pic_width, pic_height)
                hist_image(prs, slide, width, height, g, image_on_slide)
                image_crop(prs, slide, width, height, image_on_slide, g)
                image_on_slide = False
                inserted_images += 1
            if (len(files)) <= inserted_images:
                slide_counter = len(files) - 1
        slide_counter += 1
        image_on_slide = True
        attributeOn = set_image_stats(g, obj)

        merged_list = first_stats_list[0] + second_stats_list[0]
        item_list = first_stats_list[1]
        # print (item_list)
        create_table(prs, slide, item_list, merged_list, device1_name, device2_name)
        if attributeOn[2] == "1":
            obj.set_attribute_to_slide(prs, slide)
        first_stats_list[0].clear()
        second_stats_list[0].clear()
        if slide_counter == len(files):
            break


# Create table
def create_table(prs, slide, table_default, stats_list, device1_name, device2_name):
    table_left = int(prs.slide_width * 0.76)
    table_top = int(prs.slide_height * 0.08)
    table_width = int(prs.slide_width / 3)
    table_height = int(prs.slide_height / 2.0)

    count = 0
    shape = slide.shapes.add_table(int(len(stats_list) / 2) + 1, 3, table_left, table_top,
                                   table_width, table_height)
    table = shape.table
    table.columns[0].width = Inches(1.4)
    table.columns[1].width = Inches(0.8)
    table.columns[2].width = Inches(0.8)
    for row in range(int(len(stats_list) / 2) + 1):
        for cols in range(3):
            fill = table.cell(row, cols).fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(229, 229, 229)
    for row in range(int(len(stats_list) / 2) + 1):
        for cols in range(1):
            cell = table.cell(row, cols)
            l = table.cell(row, cols).text_frame.paragraphs[0]
            run = l.add_run()
            run.text = table_default[count]
            font = run.font
            font.size = Pt(12)
            font.color.rgb = RGBColor(0, 0, 0)
            count += 1
    count = 0

    cell = table.cell(0, 2)
    l = table.cell(0, 2).text_frame.paragraphs[0]
    run = l.add_run()
    run.text = device2_name
    font = run.font
    font.size = Pt(12)
    font.color.rgb = RGBColor(0, 0, 0)

    cell = table.cell(0, 1)
    l = table.cell(0, 1).text_frame.paragraphs[0]
    run = l.add_run()
    run.text = device1_name
    font = run.font
    font.size = Pt(12)
    font.color.rgb = RGBColor(0, 0, 0)

    for row in range(1, int(len(stats_list) / 2) + 1):
        cell = table.cell(row, 1)
        l = table.cell(row, 1).text_frame.paragraphs[0]
        run = l.add_run()
        run.text = stats_list[count]
        font = run.font
        font.size = Pt(12)
        count += 1
    if len(stats_list) == len(stats_list):
        for row in range(1, int(len(stats_list) / 2) + 1):
            cell = table.cell(row, 2)
            l = table.cell(row, 2).text_frame.paragraphs[0]
            run = l.add_run()
            run.text = stats_list[count]
            font = run.font
            font.size = Pt(12)
            count += 1
            ten_list = [10]
            five_list = [2]
            zero_list = [0]
            if count < 28:
                first_value = re.findall(r'\d+\.\d+|\d+', stats_list[count - 10])
                second_value = re.findall(r'\d+\.\d+|\d+', stats_list[count - 1])
            elif count == 28:
                first_value = re.findall(r'%s(\d+)' % '/', stats_list[count - 10])
                second_value = re.findall(r'%s(\d+)' % '/', stats_list[count - 1])
            first_value = list(map(float, first_value))
            second_value = list(map(float, second_value))
            ten_pro = list(map(truediv, first_value, ten_list))
            five_pro = list(map(truediv, ten_pro, five_list))
            subtraction_stats = list(map(float.__sub__, first_value, second_value))
            if subtraction_stats > zero_list:
                negative_ten_pro = list(map(float.__sub__, ten_pro, subtraction_stats))
                negative_five_pro = list(map(float.__sub__, five_pro, subtraction_stats))
            else:
                negative_ten_pro = list(map(float.__add__, ten_pro, subtraction_stats))
                negative_five_pro = list(map(float.__add__, five_pro, subtraction_stats))
            if ten_pro < subtraction_stats:
                fill = table.cell(row, 2).fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(229, 229, 229)
            elif ten_pro > subtraction_stats > five_pro:
                fill = table.cell(row, 2).fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(229, 229, 229)
            elif subtraction_stats < five_pro and negative_five_pro > zero_list:
                fill = table.cell(row, 2).fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(229, 229, 229)
            elif negative_ten_pro > zero_list > negative_five_pro:
                fill = table.cell(row, 2).fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(229, 229, 229)
            elif subtraction_stats < zero_list and negative_ten_pro < zero_list:
                fill = table.cell(row, 2).fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(229, 229, 229)
            subtraction_stats = []
            ten_pro = []

    count = 0


# Image to crop
def image_crop(prs, slide, width, height, image_on_slide, image_to_crop):
    if height < width:
        crop_left = int(prs.slide_width * 0.01)
        crop_top = int(prs.slide_height * 0.53)
        crop_width = int(prs.slide_width / 5)
        crop_height = int(crop_width * height / width)
        if image_on_slide == False:
            if (height / width) > 0.65:
                crop_left += int(prs.slide_width * 0.34)
                crop_top = int(prs.slide_height * 0.53)
            else:
                crop_left += int(prs.slide_width * 0.37)  # 16:9
                crop_top = int(prs.slide_height * 0.50)
        else:
            if (height / width) > 0.65:
                crop_left = int(prs.slide_width * 0.144)
                crop_top = int(prs.slide_height * 0.53)
            else:
                crop_left = int(prs.slide_width * 0.172)  # 16:9
                crop_top = int(prs.slide_height * 0.50)
    else:
        crop_width = int(prs.slide_width / 7.3)
        crop_height = int(crop_width / (width / height))
        crop_left = int(prs.slide_width * 0.01)
        crop_top = int(prs.slide_height * 0.08)
        if not image_on_slide:
            crop_left += int(prs.slide_width * 0.535)
        else:
            crop_left = int(prs.slide_width * 0.4)
    pic = slide.shapes.add_picture(image_to_crop, crop_left, crop_top,
                                   crop_width, crop_height)


# Historgram dump
def hist_image(prs, slide, width, height, image_hist, image_on_slide):
    imgz = cv2.imread(image_hist, 0)
    plt.hist(imgz.ravel(), 256, [0, 256], color="black");
    plt.plot()
    plt.xticks([])
    plt.yticks([])
    plt.savefig('histogram.png')
    img_histogram = Image.open('histogram.png')
    left = 100
    top = 60
    right = 565
    bottom = 430
    img_histogram = img_histogram.crop((left, top, right, bottom))
    img_histogram = img_histogram.save('histogram.png')

    # if img.shape[0] < img.shape[1]:
    if height < width:
        hist_left = int(prs.slide_width * 0.01)
        hist_top = int(prs.slide_height * 0.53)
        hist_width = int(prs.slide_width / 8)
        hist_height = int(hist_width * 9 / 16)
        if not image_on_slide:
            if (height / width) > 0.65:
                hist_left += int(prs.slide_width * 0.545)
            else:
                hist_left += int(prs.slide_width * 0.58)  # 16:9
                hist_top = int(prs.slide_height * 0.50)
        else:
            if (height / width) > 0.65:
                hist_left = int(prs.slide_width * 0.01)
            else:
                hist_left = int(prs.slide_width * 0.04)  # 16:9
                hist_top = int(prs.slide_height * 0.50)
    else:
        hist_left = int(prs.slide_width * 0.4)
        hist_top = int(prs.slide_height * 0.42)
        hist_width = int(prs.slide_width / 8)
        hist_height = int(hist_width * 9 / 16)
        if not image_on_slide:
            hist_left += int(prs.slide_width * 0.145)
        else:
            hist_left = int(prs.slide_width * 0.4)

    pic = slide.shapes.add_picture('histogram.png', hist_left, hist_top,
                                   hist_width, hist_height)
    plt.clf()
    os.remove('histogram.png')


# Get device name
def get_device_name(image_name):
    devices = image_name.split("_")[-1].split(".")[0]
    devices = devices[1:-1]
    return devices


# Set device name
def set_device_name(prs, image_name, slide, image_on_slide, width, height):
    name_device = get_device_name(image_name)
    if height < width:
        name_left = int(prs.slide_width * 0.01)
        name_top = int(prs.slide_height * 0.68)
        name_width = int(prs.slide_width / 8)
        name_height = int(name_width * 9 / 16)
        if not image_on_slide:
            if (height / width) > 0.65:
                name_left += int(prs.slide_width * 0.545)
            else:
                name_left += int(prs.slide_width * 0.61)
                name_top = int(prs.slide_height * 0.65)  # 16:9
        else:
            if (height / width) > 0.65:
                name_left = int(prs.slide_width * 0.01)
            else:
                name_left = int(prs.slide_width * 0.04)
                name_top = int(prs.slide_height * 0.65)  # 16:9
    else:
        name_left = int(prs.slide_width * 0.4)
        name_top = int(prs.slide_height * 0.55)
        name_width = int(prs.slide_width / 8)
        name_height = int(name_width * 9 / 16)
        if not image_on_slide:
            name_left += int(prs.slide_width * 0.145)
        else:
            name_left = int(prs.slide_width * 0.4)

    txBox = slide.shapes.add_textbox(name_left, name_top,
                                     name_width, name_height)
    tf = txBox.text_frame
    tf.text = name_device
    tf.paragraphs[0].font.size = Pt(12)


# Set title to a slide
def set_title_to_slide(image_name, title, g, prs):
    slide_title = os.path.basename(os.path.dirname(g))
    case_num = image_name.split("_")[0]
    case_num = case_num[0:-1]

    title.text = case_num + " : " + slide_title
    title.text_frame.paragraphs[0].font.size = Pt(24)
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    return slide_title


# Sharpness
def sharpness_image(img):
    array = np.asarray(img, dtype=np.int32)

    gy, gx = np.gradient(array)
    gnorm = np.sqrt(gx ** 2 + gy ** 2)
    img_sharpness = np.average(gnorm)

    return int(round(img_sharpness))


# Dynamic Range
def dynamic_range(img):
    img_gray = np.maximum(img[:, :, 0], img[:, :, 1], img[:, :, 2])

    hist, bins = np.histogram(img_gray, 256, [0, 256])
    cdf = hist.cumsum()
    cdf_nom = cdf / cdf.max()
    plot(cdf_nom)
    over_exposed = int(100 - cdf_nom[204] * 100)
    under_exposed = int(cdf_nom[3] * 100)
    dynamic_range = 100 - over_exposed - under_exposed

    return [over_exposed, under_exposed, dynamic_range]


# Peak saturation
def peak_saturation(img):
    rgb2hsv = np.vectorize(colorsys.rgb_to_hsv)

    img_float = img / 255.0
    h, s, v = rgb2hsv(img_float[:, :, 0], img_float[:, :, 1], img_float[:, :, 2])

    hist, bins = np.histogram(s, 100, [0.0, 1.0])
    peaks, properties = find_peaks(hist, prominence=1, distance=5)

    if len(peaks) > 0:
        peak1_sat = peaks[0]
    else:
        peak1_sat = None

    if len(peaks) > 1:
        peak_inx = np.argmax(hist[peaks[1:]])
        peak2_sat = peaks[1 + peak_inx]
    else:
        peak2_sat = None

    hist, bins = np.histogram(h, 100, [0.0, 1.0])
    peaks, properties = find_peaks(hist, prominence=1, distance=10)

    if len(peaks) > 0:
        peak1_hue = int(peaks[0] / 100 * 360 + 0.5)
    else:
        peak1_hue = None
    if len(peaks) > 1:
        peak_inx = np.argmax(hist[peaks[1:]])
        peak2_hue = int(peaks[1 + peak_inx] / 100 * 360 + 0.5)
    else:
        peak2_hue = None

    return [peak1_sat, peak1_hue, peak2_sat, peak2_hue]


# Set image stats
def set_image_stats(image_name, obj):
    img = Image.open(image_name)
    rgb_array = np.asarray(img)
    rgb_copy = rgb_array.copy()
    stats_list = obj.show_req(rgb_copy, image_name)
    return stats_list
